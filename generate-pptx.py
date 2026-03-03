#!/usr/bin/env python3
"""Assemble slide screenshots into a PowerPoint deck."""

import glob
import os
import re
from pptx import Presentation
from pptx.util import Inches, Emu

ROOT = os.path.dirname(os.path.abspath(__file__))
SCREENSHOT_DIR = os.path.join(ROOT, "screenshots")
OUTPUT_PATH = os.path.join(ROOT, "public", "Source-AI-Pitch-Deck.pptx")

def sort_key(path: str):
    name = os.path.basename(path)
    match = re.search(r"slide-(\d+)\.png$", name)
    if match:
        return int(match.group(1))
    return name


screenshots = sorted(glob.glob(os.path.join(SCREENSHOT_DIR, "slide-*.png")), key=sort_key)

if not screenshots:
    print("ERROR: No screenshots found. Run capture-slides.mjs first.")
    exit(1)

print(f"Building PPTX from {len(screenshots)} screenshots...")

prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

for img_path in screenshots:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(img_path, Emu(0), Emu(0), prs.slide_width, prs.slide_height)
    print(f"  Added {os.path.basename(img_path)}")

prs.save(OUTPUT_PATH)
print(f"Saved to {OUTPUT_PATH}")
